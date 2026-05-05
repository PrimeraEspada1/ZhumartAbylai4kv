from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt


BASE_DIR = Path(__file__).resolve().parent
TEMPLATE = BASE_DIR / "Шаблон презентации проекта для лаб 5.pptx"
OUTPUT = BASE_DIR / "RepairHub_Lab5_Жумарт_Абылай_v2.pptx"


SLIDES = [
    [
        "RepairHub – компьютер жөндеу сервисі",
        "Студент Жумарт Абылай",
        "Мамандығы",
        "Computer Science",
        "/ Білім беру бағдарламасы: 6В06102",
        "Алматы 2026",
    ],
    [
        "Жобаның мақсаты",
        "RepairHub жобасының мақсаты –",
        "компьютер жөндеу қызметін онлайн басқару.",
        "Пайдаланушы:",
        "өтінім қалдыра алады",
        "профиль аша алады",
        "AI чат-боттан кеңес алады",
    ],
    [
        "Мақсаттар мен міндеттер",
        "",
        "Жоба міндеттері:",
        "Негізгі функциялар:",
        "тіркелу және жүйеге кіру",
        "жөндеу өтінімдерін базаға сақтау",
        "UX/UI ыңғайлылығын жақсарту",
    ],
    [
        "Деректер базасы (Database)",
        "Жобада",
        "ServiceRequest",
        "моделі",
        "",
        "жасалды",
        ":",
        "title –",
        "өтінім",
        "",
        "атауы",
        "",
        "estimated_price –",
        "бағасы",
        "",
        "status –",
        "мәртебесі",
        "",
        "Барлық",
        "",
        "деректер",
        "",
        "базаға",
        "",
        "сақталады",
    ],
    [
        "Backend жұмысы (Django)",
        "Backend",
        "арқылы",
        ":",
        "Өтінімдер",
        "",
        "өңделеді",
        "",
        "Базаға",
        "",
        "жазылады",
        "",
        "Сайтқа",
        "",
        "жіберіледі",
        "",
        "Django ORM",
        "қолданылды",
    ],
    [
        "Жоба",
        "",
        "тақырыбы",
        "",
        "бойынша",
        "",
        "шолу",
        "Жобада",
        "",
        "қолданылды",
        ":",
        "Python",
        "Django framework",
        "HTML",
        "CSS",
        "Jinja",
        "шаблон",
        "",
        "жүйесі",
    ],
    [
        "Frontend (Jinja + HTML)",
        "Frontend",
        "бөлігі",
        ":",
        "HTML",
        "арқылы",
        "",
        "құрылды",
        "",
        "Jinja",
        "арқылы",
        "",
        "динамикалық",
        "",
        "болды",
        "",
        "Өтінімдер",
        "",
        "сайтта",
        "",
        "автоматты",
        "",
        "түрде",
        "",
        "шығады",
    ],
    [
        "UX/UI талдау",
        "Тестілеу:",
        "функционалдық",
        "және юзабилити",
        "",
        "нәтижесі:",
        "форма тексерілді",
        "",
        "навигация түсінікті",
        "",
        "қателер",
        "анық",
        "көрсетілуі",
        "керек",
        "Dark / Light",
        "режим",
        "қосылды",
    ],
    [
        "UX/UI жақсарту жолдары",
        "Inline-валидация қосу",
        "FAQ және подсказка қосу",
        "Қате өрісті белгілеу",
        "Минималистік дизайн сақтау",
    ],
    [
        "Қортынды",
        "",
        "RepairHub –",
        "толық",
        "",
        "веб",
        "",
        "жоба",
        ".",
        "Бұл",
        "",
        "жоба",
        "",
        "арқылы",
        ":",
        "Django",
        "үйрендік",
        "",
        "База",
        "",
        "мен",
        "",
        "жұмыс",
        "",
        "жасадық",
        "",
        "UX/UI",
        "түсіндік",
        "",
    ],
]


def set_run_text(shape, values):
    runs = []
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            runs.append(run)

    for index, run in enumerate(runs):
        if index < len(values):
            run.text = values[index]
        else:
            run.text = ""


def normalize_text_boxes(slide):
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        shape.text_frame.word_wrap = True
        shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.size and run.font.size.pt > 30:
                    run.font.size = Pt(28)


def main():
    prs = Presentation(TEMPLATE)

    for slide, values in zip(prs.slides, SLIDES):
        text_shapes = [shape for shape in slide.shapes if getattr(shape, "has_text_frame", False)]
        runs_total = 0
        for shape in text_shapes:
            for paragraph in shape.text_frame.paragraphs:
                runs_total += len(paragraph.runs)

        cursor = 0
        for shape in text_shapes:
            shape_runs = sum(len(paragraph.runs) for paragraph in shape.text_frame.paragraphs)
            set_run_text(shape, values[cursor : cursor + shape_runs])
            cursor += shape_runs

        normalize_text_boxes(slide)

    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    main()
